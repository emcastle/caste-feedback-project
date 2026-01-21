from __future__ import annotations

import hashlib
import io
import json
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import pandas as pd
from pptx import Presentation

try:
    from PIL import Image
except Exception:
    Image = None

try:
    import pytesseract
except Exception:
    pytesseract = None


@dataclass(frozen=True)
class PptxReadConfig:
    include_speaker_notes: bool = True
    max_chars_per_block: int = 20000

    # Option A: OCR is available, but OFF by default
    ocr_images: bool = False
    ocr_lang: str = "eng"
    ocr_max_chars: int = 20000


def sha256_file(path: Path, chunk_size: int = 1024 * 1024) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        while True:
            b = f.read(chunk_size)
            if not b:
                break
            h.update(b)
    return h.hexdigest()


def _truncate(s: Optional[str], max_chars: int) -> Optional[str]:
    if s is None:
        return None
    if len(s) <= max_chars:
        return s
    return s[: max_chars - 20] + " ...[TRUNCATED]"


def _table_to_json(table) -> str:
    rows = []
    for r in table.rows:
        row_vals = []
        for c in r.cells:
            row_vals.append((c.text or "").strip())
        rows.append(row_vals)
    return json.dumps(rows, ensure_ascii=False)


def _ocr_image_bytes(image_bytes: bytes, cfg: PptxReadConfig) -> Tuple[Optional[str], Optional[str]]:
    """
    Returns (ocr_text, error). Never raises.
    """
    if not cfg.ocr_images:
        return None, None

    if Image is None:
        return None, "PIL not available; cannot OCR images"

    if pytesseract is None:
        return None, "pytesseract not available; cannot OCR images"

    try:
        img = Image.open(io.BytesIO(image_bytes))
        text = pytesseract.image_to_string(img, lang=cfg.ocr_lang).strip()
        text = _truncate(text, cfg.ocr_max_chars)
        return text, None
    except Exception as e:
        return None, f"{type(e).__name__}: {e}"


def extract_pptx_to_relational(
    pptx_path: Path,
    source_rel_path: str,
    cfg: PptxReadConfig = PptxReadConfig(),
) -> Tuple[pd.DataFrame, pd.DataFrame]:
    """
    Extract PPTX into:
      - documents_df: 1 row per file
      - blocks_df: ordered blocks across slides

    blocks_df columns include image OCR outputs when enabled:
      - ocr_text, ocr_text_len, image_content_type
    """
    pptx_path = pptx_path.resolve()
    doc_id = sha256_file(pptx_path)

    source_file = pptx_path.name
    source_ext = ".pptx"
    source_type = "pptx"
    extractor_used = "python-pptx"

    try:
        prs = Presentation(str(pptx_path))
        blocks: List[Dict[str, Any]] = []

        slide_count = len(prs.slides)
        num_images = 0
        num_images_ocr_success = 0
        num_images_ocr_failed = 0
        num_tables = 0

        block_num_global = 0

        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                # Table
                if hasattr(shape, "has_table") and shape.has_table:
                    num_tables += 1
                    table_json = _table_to_json(shape.table)
                    blocks.append(
                        {
                            "doc_id": doc_id,
                            "slide_num": slide_idx,
                            "block_num": block_num_global,
                            "block_type": "table",
                            "shape_num": shape_idx,
                            "text": _truncate(table_json, cfg.max_chars_per_block),
                            "has_image": False,
                            "image_content_type": None,
                            "ocr_text": None,
                            "ocr_text_len": 0,
                            "meta_json": json.dumps(
                                {"rows": len(shape.table.rows), "cols": len(shape.table.columns)},
                                ensure_ascii=False,
                            ),
                            "error": None,
                        }
                    )
                    block_num_global += 1
                    continue

                # Text box / placeholder text
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    t = (shape.text_frame.text or "").strip()
                    if t:
                        blocks.append(
                            {
                                "doc_id": doc_id,
                                "slide_num": slide_idx,
                                "block_num": block_num_global,
                                "block_type": "text",
                                "shape_num": shape_idx,
                                "text": _truncate(t, cfg.max_chars_per_block),
                                "has_image": False,
                                "image_content_type": None,
                                "ocr_text": None,
                                "ocr_text_len": 0,
                                "meta_json": None,
                                "error": None,
                            }
                        )
                        block_num_global += 1
                    continue

                # Images (pictures)
                if hasattr(shape, "image"):
                    num_images += 1
                    image_bytes = None
                    image_ct = None
                    try:
                        image_bytes = shape.image.blob
                        image_ct = getattr(shape.image, "content_type", None)
                    except Exception:
                        image_bytes = None

                    ocr_text = None
                    ocr_err = None
                    if image_bytes:
                        ocr_text, ocr_err = _ocr_image_bytes(image_bytes, cfg)
                        if cfg.ocr_images:
                            if ocr_err is None:
                                num_images_ocr_success += 1
                            else:
                                num_images_ocr_failed += 1
                    else:
                        if cfg.ocr_images:
                            num_images_ocr_failed += 1
                        ocr_err = "Could not extract image bytes from PPTX shape"

                    blocks.append(
                        {
                            "doc_id": doc_id,
                            "slide_num": slide_idx,
                            "block_num": block_num_global,
                            "block_type": "image",
                            "shape_num": shape_idx,
                            "text": None,
                            "has_image": True,
                            "image_content_type": image_ct,
                            "ocr_text": ocr_text,
                            "ocr_text_len": len(ocr_text) if ocr_text else 0,
                            "meta_json": None,
                            "error": ocr_err,
                        }
                    )
                    block_num_global += 1
                    continue

            # Speaker notes
            if cfg.include_speaker_notes:
                try:
                    notes_slide = slide.notes_slide
                    notes_text = (notes_slide.notes_text_frame.text or "").strip()
                    if notes_text:
                        blocks.append(
                            {
                                "doc_id": doc_id,
                                "slide_num": slide_idx,
                                "block_num": block_num_global,
                                "block_type": "notes",
                                "shape_num": None,
                                "text": _truncate(notes_text, cfg.max_chars_per_block),
                                "has_image": False,
                                "image_content_type": None,
                                "ocr_text": None,
                                "ocr_text_len": 0,
                                "meta_json": None,
                                "error": None,
                            }
                        )
                        block_num_global += 1
                except Exception:
                    pass

        blocks_df = pd.DataFrame(blocks)

        documents_df = pd.DataFrame(
            [
                {
                    "doc_id": doc_id,
                    "source_file": source_file,
                    "source_rel_path": source_rel_path,
                    "source_ext": source_ext,
                    "source_type": source_type,
                    "num_slides": int(slide_count),
                    "num_blocks": int(len(blocks_df)),
                    "num_images": int(num_images),
                    "num_images_ocr_success": int(num_images_ocr_success),
                    "num_images_ocr_failed": int(num_images_ocr_failed),
                    "num_tables": int(num_tables),
                    "extractor_used": extractor_used,
                    "error": None,
                }
            ]
        )

        return documents_df, blocks_df

    except Exception as e:
        documents_df = pd.DataFrame(
            [
                {
                    "doc_id": doc_id,
                    "source_file": source_file,
                    "source_rel_path": source_rel_path,
                    "source_ext": source_ext,
                    "source_type": source_type,
                    "num_slides": None,
                    "num_blocks": None,
                    "num_images": None,
                    "num_images_ocr_success": None,
                    "num_images_ocr_failed": None,
                    "num_tables": None,
                    "extractor_used": extractor_used,
                    "error": f"{type(e).__name__}: {e}",
                }
            ]
        )

        blocks_df = pd.DataFrame(
            [
                {
                    "doc_id": doc_id,
                    "slide_num": None,
                    "block_num": None,
                    "block_type": None,
                    "shape_num": None,
                    "text": None,
                    "has_image": None,
                    "image_content_type": None,
                    "ocr_text": None,
                    "ocr_text_len": 0,
                    "meta_json": None,
                    "error": f"{type(e).__name__}: {e}",
                }
            ]
        )
        return documents_df, blocks_df
